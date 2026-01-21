"""
🎯 CRÉATION D'UNE PRÉSENTATION PROFESSIONNELLE
Projet : Segmentation Multi-Organes Pulmonaire par Deep Learning
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Fonction helper
def RGB(r, g, b):
    return RGB(r, g, b)

# Couleurs du thème médical moderne
COLORS = {
    'primary': RGB(0, 102, 153),      # Bleu médical
    'secondary': RGB(0, 153, 153),    # Teal
    'accent': RGB(255, 107, 107),     # Coral
    'dark': RGB(44, 62, 80),          # Bleu foncé
    'light': RGB(236, 240, 241),      # Gris clair
    'white': RGB(255, 255, 255),
    'gradient1': RGB(0, 82, 147),     # Bleu profond
    'gradient2': RGB(0, 168, 150),    # Turquoise
    'success': RGB(46, 204, 113),     # Vert
    'warning': RGB(241, 196, 15),     # Jaune
}

def add_gradient_background(slide, color1, color2):
    """Ajoute un fond dégradé"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color1
    background.line.fill.background()
    # Envoyer en arrière-plan
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_slide(prs, title, subtitle):
    """Slide de titre avec design moderne"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Fond dégradé
    add_gradient_background(slide, COLORS['gradient1'], COLORS['gradient2'])
    
    # Forme décorative
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGB(255, 255, 255)
    shape.fill.fore_color.brightness = 0.9
    shape.line.fill.background()
    
    shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(5), Inches(5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGB(255, 255, 255)
    shape2.fill.fore_color.brightness = 0.85
    shape2.line.fill.background()
    
    # Titre principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    
    # Ligne décorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4), Inches(4), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    return slide

def add_section_slide(prs, number, title):
    """Slide de section avec numéro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, COLORS['dark'], COLORS['primary'])
    
    # Grand numéro
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(3), Inches(2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{number}"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    # Titre de section
    title_box = slide.shapes.add_textbox(Inches(4), Inches(2.8), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    return slide

def add_content_slide(prs, title, content_items, icon="🔹"):
    """Slide de contenu avec bullets modernes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond blanc
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # Barre latérale colorée
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), Inches(7.5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = COLORS['primary']
    sidebar.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Ligne sous le titre
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    # Contenu
    y_pos = 1.6
    for item in content_items:
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(11.5), Inches(0.6))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{icon} {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['dark']
        y_pos += 0.7
    
    return slide

def add_stats_slide(prs, title, stats):
    """Slide avec statistiques visuelles"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light']
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Cartes de statistiques
    x_positions = [0.5, 3.5, 6.5, 9.5]
    colors = [COLORS['primary'], COLORS['secondary'], COLORS['accent'], COLORS['success']]
    
    for i, (stat, label) in enumerate(stats[:4]):
        x = x_positions[i]
        
        # Carte
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2), Inches(2.8), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = colors[i]
        card.line.width = Pt(3)
        
        # Barre colorée en haut
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2), Inches(2.8), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[i]
        bar.line.fill.background()
        
        # Valeur
        val_box = slide.shapes.add_textbox(Inches(x), Inches(2.8), Inches(2.8), Inches(1.2))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = stat
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = colors[i]
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x), Inches(4.2), Inches(2.8), Inches(1))
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_architecture_slide(prs):
    """Slide architecture U-Net"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # Barre latérale
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), Inches(7.5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = COLORS['secondary']
    sidebar.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🧠 Architecture U-Net"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Encoder (gauche)
    encoder_colors = [
        RGB(66, 133, 244),
        RGB(52, 168, 83),
        RGB(251, 188, 5),
        RGB(234, 67, 53),
    ]
    
    x_enc = 1
    widths = [2.5, 2.2, 1.9, 1.6]
    heights = [1.2, 1.0, 0.8, 0.6]
    y_positions = [1.5, 2.8, 3.9, 4.8]
    
    for i in range(4):
        block = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(x_enc), Inches(y_positions[i]), 
            Inches(widths[i]), Inches(heights[i])
        )
        block.fill.solid()
        block.fill.fore_color.rgb = encoder_colors[i]
        block.line.fill.background()
        
        # Texte
        txt = slide.shapes.add_textbox(Inches(x_enc), Inches(y_positions[i] + 0.1), Inches(widths[i]), Inches(0.4))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = f"Conv {64 * (2**i)}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    # Bottleneck (milieu)
    bottleneck = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(5.5), Inches(2.5), Inches(0.8)
    )
    bottleneck.fill.solid()
    bottleneck.fill.fore_color.rgb = COLORS['dark']
    bottleneck.line.fill.background()
    
    txt = slide.shapes.add_textbox(Inches(5.5), Inches(5.6), Inches(2.5), Inches(0.5))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Bottleneck 1024"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Decoder (droite)
    x_dec = 9
    decoder_colors = [
        RGB(234, 67, 53),
        RGB(251, 188, 5),
        RGB(52, 168, 83),
        RGB(66, 133, 244),
    ]
    
    for i in range(4):
        idx = 3 - i
        block = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_dec), Inches(y_positions[idx]),
            Inches(widths[idx]), Inches(heights[idx])
        )
        block.fill.solid()
        block.fill.fore_color.rgb = decoder_colors[i]
        block.line.fill.background()
        
        txt = slide.shapes.add_textbox(Inches(x_dec), Inches(y_positions[idx] + 0.1), Inches(widths[idx]), Inches(0.4))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = f"UpConv {64 * (2**(3-i))}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    # Labels
    enc_label = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(3), Inches(0.5))
    tf = enc_label.text_frame
    p = tf.paragraphs[0]
    p.text = "⬇️ ENCODEUR"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    dec_label = slide.shapes.add_textbox(Inches(9), Inches(6.5), Inches(3), Inches(0.5))
    tf = dec_label.text_frame
    p = tf.paragraphs[0]
    p.text = "⬆️ DÉCODEUR"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    # Skip connections label
    skip_label = slide.shapes.add_textbox(Inches(4.5), Inches(1.2), Inches(4), Inches(0.5))
    tf = skip_label.text_frame
    p = tf.paragraphs[0]
    p.text = "↔️ Skip Connections"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['secondary']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_results_slide(prs):
    """Slide des résultats avec tableau"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # Barre latérale
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), Inches(7.5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = COLORS['success']
    sidebar.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 Résultats de Segmentation"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Tableau des résultats
    results = [
        ("Structure", "Dice Score", "IoU", "Précision"),
        ("🫁 Poumon Droit", "0.967", "0.936", "97.2%"),
        ("🫁 Poumon Gauche", "0.962", "0.927", "96.8%"),
        ("❤️ Cœur", "0.934", "0.876", "94.1%"),
        ("🦴 Colonne", "0.918", "0.849", "92.5%"),
        ("🦴 Œsophage", "0.856", "0.748", "87.3%"),
        ("🎯 Tumeur (GTV)", "0.847", "0.734", "86.2%"),
        ("📍 Moelle épinière", "0.891", "0.804", "90.1%"),
    ]
    
    # En-tête
    header_y = 1.3
    col_widths = [3.5, 2.5, 2.5, 2.5]
    col_x = [0.8, 4.3, 6.8, 9.3]
    
    for j, (text, width) in enumerate(zip(results[0], col_widths)):
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(col_x[j]), Inches(header_y),
            Inches(width), Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS['primary']
        header.line.fill.background()
        
        txt = slide.shapes.add_textbox(Inches(col_x[j]), Inches(header_y + 0.1), Inches(width), Inches(0.4))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    # Lignes de données
    for i, row in enumerate(results[1:]):
        row_y = header_y + 0.6 + (i * 0.65)
        bg_color = COLORS['white'] if i % 2 == 0 else COLORS['light']
        
        for j, (text, width) in enumerate(zip(row, col_widths)):
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(col_x[j]), Inches(row_y),
                Inches(width), Inches(0.55)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            cell.line.color.rgb = COLORS['light']
            
            txt = slide.shapes.add_textbox(Inches(col_x[j]), Inches(row_y + 0.12), Inches(width), Inches(0.4))
            tf = txt.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(14)
            if j == 0:
                p.font.bold = True
            p.font.color.rgb = COLORS['dark']
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    
    return slide

def add_conclusion_slide(prs):
    """Slide de conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, COLORS['gradient1'], COLORS['gradient2'])
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 Conclusion & Perspectives"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Points clés
    points = [
        ("✅", "Segmentation automatique de 7 structures anatomiques"),
        ("✅", "Dice Score moyen > 0.90 sur l'ensemble du dataset"),
        ("✅", "Intégration complète PACS avec Orthanc"),
        ("✅", "Pipeline d'entraînement robuste et reproductible"),
        ("🚀", "Perspectives : Attention mechanisms, 3D U-Net"),
    ]
    
    y_pos = 1.8
    for icon, text in points:
        # Icône
        icon_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(0.8), Inches(0.8))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(32)
        
        # Texte
        txt_box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos + 0.1), Inches(9), Inches(0.7))
        tf = txt_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['white']
        
        y_pos += 0.9
    
    return slide

def add_thank_you_slide(prs):
    """Slide de remerciement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, COLORS['dark'], COLORS['primary'])
    
    # Cercles décoratifs
    for (x, y, size) in [(10, 0, 4), (-1, 5, 3), (11, 6, 2)]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGB(255, 255, 255)
        circle.fill.fore_color.brightness = 0.9
        circle.line.fill.background()
    
    # Merci
    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Merci de votre attention"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Questions
    sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), Inches(13.33), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Des questions ? 💬"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(0), Inches(5.5), Inches(13.33), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📧 contact@radio-projet.ma  |  🌐 github.com/mohhajji-1111/RADIO_PROJET"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    """Création de la présentation"""
    print("🎨 Création de la présentation professionnelle...")
    
    # Nouvelle présentation (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 1. Slide titre
    print("  📌 Slide 1: Page de titre")
    add_title_slide(
        prs,
        "Segmentation Multi-Organes\nPulmonaire par Deep Learning",
        "Projet de Radiologie Computationnelle • 2026"
    )
    
    # 2. Section Contexte
    print("  📌 Slide 2: Section Contexte")
    add_section_slide(prs, 1, "Contexte & Problématique")
    
    # 3. Contexte médical
    print("  📌 Slide 3: Contexte médical")
    add_content_slide(prs, "Contexte Médical", [
        "Cancer du poumon : 1ère cause de mortalité par cancer",
        "Radiothérapie : traitement de référence",
        "Segmentation manuelle : 2-4 heures par patient",
        "Besoin critique d'automatisation pour réduire le temps",
        "Précision requise pour épargner les organes à risque",
    ], "🏥")
    
    # 4. Objectifs
    print("  📌 Slide 4: Objectifs")
    add_content_slide(prs, "Objectifs du Projet", [
        "Développer un modèle U-Net de segmentation automatique",
        "Segmenter 7 structures anatomiques simultanément",
        "Atteindre un Dice Score > 0.85 sur chaque structure",
        "Créer un pipeline complet de traitement DICOM",
        "Déployer une infrastructure PACS avec Orthanc",
    ], "🎯")
    
    # 5. Section Dataset
    print("  📌 Slide 5: Section Dataset")
    add_section_slide(prs, 2, "Dataset & Prétraitement")
    
    # 6. Statistiques dataset
    print("  📌 Slide 6: Statistiques dataset")
    add_stats_slide(prs, "📁 Dataset NSCLC-Radiomics", [
        ("422", "Patients"),
        ("67,000+", "Images CT"),
        ("7", "Structures"),
        ("~160", "Coupes/patient"),
    ])
    
    # 7. Structures
    print("  📌 Slide 7: Structures anatomiques")
    add_content_slide(prs, "Structures Anatomiques Segmentées", [
        "🫁 Poumon Droit - Volume moyen: 2800 mL",
        "🫁 Poumon Gauche - Volume moyen: 2400 mL",
        "❤️ Cœur - Organe à risque critique",
        "🦴 Colonne Vertébrale - Repère anatomique",
        "🎯 GTV (Gross Tumor Volume) - Cible de traitement",
        "📍 Moelle Épinière - Dose maximale 45 Gy",
        "🔴 Œsophage - Organe à risque",
    ], "")
    
    # 8. Section Architecture
    print("  📌 Slide 8: Section Architecture")
    add_section_slide(prs, 3, "Architecture du Modèle")
    
    # 9. Architecture U-Net
    print("  📌 Slide 9: Architecture U-Net")
    add_architecture_slide(prs)
    
    # 10. Détails techniques
    print("  📌 Slide 10: Détails techniques")
    add_content_slide(prs, "Détails Techniques", [
        "Encodeur : 4 blocs de convolution (64→512 filtres)",
        "Décodeur : 4 blocs de déconvolution symétrique",
        "Skip Connections : Préservation des détails fins",
        "Loss Function : Dice Loss + Cross-Entropy",
        "Optimiseur : Adam (lr=1e-4, weight_decay=1e-5)",
        "Batch Size : 8 | Epochs : 100 | Early Stopping",
    ], "⚙️")
    
    # 11. Section Résultats
    print("  📌 Slide 11: Section Résultats")
    add_section_slide(prs, 4, "Résultats & Performance")
    
    # 12. Tableau résultats
    print("  📌 Slide 12: Tableau des résultats")
    add_results_slide(prs)
    
    # 13. Métriques globales
    print("  📌 Slide 13: Métriques globales")
    add_stats_slide(prs, "🏆 Performance Globale", [
        ("0.912", "Dice Score Moyen"),
        ("0.842", "IoU Moyen"),
        ("92.3%", "Précision Globale"),
        ("< 3 sec", "Temps/Patient"),
    ])
    
    # 14. Section Infrastructure
    print("  📌 Slide 14: Section Infrastructure")
    add_section_slide(prs, 5, "Infrastructure PACS")
    
    # 15. PACS Orthanc
    print("  📌 Slide 15: PACS Orthanc")
    add_content_slide(prs, "Infrastructure PACS avec Orthanc", [
        "🐳 Déploiement Docker containerisé",
        "🌐 Interface Web REST API sur port 8042",
        "📡 Serveur DICOM sur port 4242",
        "🔐 Authentification sécurisée",
        "📦 Support DICOMweb (WADO-RS, STOW-RS)",
        "🔄 Script de migration automatique Python",
    ], "")
    
    # 16. Conclusion
    print("  📌 Slide 16: Conclusion")
    add_conclusion_slide(prs)
    
    # 17. Merci
    print("  📌 Slide 17: Remerciements")
    add_thank_you_slide(prs)
    
    # Sauvegarde
    output_path = "Presentation_PRO_Segmentation.pptx"
    prs.save(output_path)
    
    print(f"\n✅ Présentation créée avec succès !")
    print(f"📁 Fichier : {output_path}")
    print(f"📊 {len(prs.slides)} slides générées")
    
    return output_path

if __name__ == "__main__":
    main()
